import os
import fitz  # PyMuPDF

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None


class FileExtractor:
    def __init__(self):
        pass

    def extract_text(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            print(f"파일을 찾을 수 없습니다: {file_path}")
            return ""

        try:
            ext = os.path.splitext(file_path)[1].lower()
            text_parts = []

            if ext == ".pdf":
                doc = fitz.open(file_path)
                total_pages = len(doc)
                
                # 큰 파일: 홀수 페이지만 읽기 (최대 15페이지)
                if total_pages > 10:
                    pages_read = 0
                    max_pages = 15
                    for i in range(0, total_pages, 2):
                        if pages_read >= max_pages:
                            break
                        page_text = doc[i].get_text()
                        if len(page_text.strip()) >= 10:
                            text_parts.append(page_text)
                            pages_read += 1
                # 작은 파일: 모든 페이지 읽기
                else:
                    for page in doc:
                        page_text = page.get_text()
                        if len(page_text.strip()) >= 10:
                            text_parts.append(page_text)
                doc.close()
            
            elif ext == ".pptx":
                if not Presentation:
                    print("python-pptx가 설치되지 않았습니다.")
                    return ""
                
                doc = Presentation(file_path)
                total_slides = len(doc.slides)
                
                # 큰 파일: 홀수 페이지만 읽기 (최대 15페이지)
                if total_slides > 10:
                    pages_read = 0
                    max_pages = 15
                    for i in range(0, total_slides, 2):
                        if pages_read >= max_pages:
                            break
                        slide = doc.slides[i]
                        slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                        if len(slide_text.strip()) >= 10:
                            text_parts.append(slide_text)
                            pages_read += 1
                # 작은 파일: 모든 페이지 읽기
                else:
                    for slide in doc.slides:
                        slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                        if len(slide_text.strip()) >= 10:
                            text_parts.append(slide_text)

            elif ext == ".docx":
                if not Document:
                    print("python-docx가 설치되지 않았습니다.")
                    return ""
                
                doc = Document(file_path)
                total_paragraphs = len(doc.paragraphs)
                
                # 큰 파일: 홀수 페이지만 읽기 (최대 15페이지)
                if total_paragraphs > 10:
                    pages_read = 0
                    max_pages = 15
                    for i in range(0, total_paragraphs, 2):
                        if pages_read >= max_pages:
                            break
                        paragraph_text = doc.paragraphs[i].text
                        if len(paragraph_text.strip()) >= 10:
                            text_parts.append(paragraph_text)
                            pages_read += 1
                # 작은 파일: 모든 페이지 읽기
                else:
                    for paragraph in doc.paragraphs:
                        paragraph_text = paragraph.text
                        if len(paragraph_text.strip()) >= 10:
                            text_parts.append(paragraph_text)

            else:
                # 지원하지 않는 형식은 파일명만 반환
                return os.path.basename(file_path)

            return "\n".join(text_parts)

        except Exception as e:
            print(f"텍스트 추출 실패 {file_path}: {e}")
            return ""

    def extract_one_page(self, file_path: str) -> str:
        """첫 번째 페이지/슬라이드의 텍스트만 추출합니다."""
        if not os.path.exists(file_path):
            print(f"파일을 찾을 수 없습니다: {file_path}")
            return ""

        try:
            ext = os.path.splitext(file_path)[1].lower()

            if ext == ".pdf":
                with fitz.open(file_path) as doc:
                    if len(doc) > 0:
                        return doc[0].get_text()
                    return ""

            elif ext == ".pptx":
                if not Presentation:
                    print("python-pptx가 설치되지 않았습니다.")
                    return ""

                prs = Presentation(file_path)
                if len(prs.slides) > 0:
                    first_slide = prs.slides[0]
                    return " ".join(
                        shape.text
                        for shape in first_slide.shapes
                        if hasattr(shape, "text")
                    )
                return ""

            elif ext == ".docx":
                if not Document:
                    print("python-docx가 설치되지 않았습니다.")
                    return ""

                doc = Document(file_path)
                if len(doc.paragraphs) > 0:
                    return doc.paragraphs[0].text
                return ""

            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.readline()

            else:
                return os.path.basename(file_path)

        except Exception as e:
            print(f"첫 페이지 추출 실패 {file_path}: {e}")
            return ""

    def extract_metadata(self, file_path):
        """파일에서 메타데이터를 추출합니다."""
        metadata = {"author": "", "title": ""}

        try:
            ext = os.path.splitext(file_path)[1].lower()

            if ext == ".pdf":
                import fitz

                with fitz.open(file_path) as doc:
                    pdf_metadata = doc.metadata
                    metadata["author"] = pdf_metadata.get("author", "")
                    metadata["title"] = pdf_metadata.get("title", "")

            elif ext == ".pptx":
                from pptx import Presentation

                prs = Presentation(file_path)
                core_props = prs.core_properties
                metadata["author"] = core_props.author or ""
                metadata["title"] = core_props.title or ""

            elif ext == ".docx":
                from docx import Document

                doc = Document(file_path)
                core_props = doc.core_properties
                metadata["author"] = doc.core_properties.author or ""
                metadata["title"] = doc.core_properties.title or ""
            
            return metadata

        except Exception as e:
            print(f"Error extracting metadata from {file_path}: {e}")
            return metadata